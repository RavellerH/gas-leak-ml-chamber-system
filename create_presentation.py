from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("GasLeak_Presentation.pptx")
CHAMBER_IMAGE = Path("Gas Test Chamber Design.jpeg")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

NAVY = RGBColor(18, 35, 52)
TEAL = RGBColor(20, 126, 129)
GREEN = RGBColor(74, 143, 89)
AMBER = RGBColor(201, 129, 45)
RED = RGBColor(166, 64, 64)
LIGHT = RGBColor(245, 248, 250)
MID = RGBColor(105, 116, 126)
WHITE = RGBColor(255, 255, 255)


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_bg(slide, color=LIGHT):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_fill(rect, color)
    slide.shapes._spTree.remove(rect._element)
    slide.shapes._spTree.insert(2, rect._element)


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.7))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(1.0), Inches(11.8), Inches(0.35))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(12)
        sp.font.color.rgb = MID


def add_footer(slide, index):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(7.05), Inches(12.25), Inches(0.02))
    set_fill(line, RGBColor(214, 222, 228))
    foot = slide.shapes.add_textbox(Inches(0.55), Inches(7.12), Inches(12.2), Inches(0.25))
    p = foot.text_frame.paragraphs[0]
    p.text = f"Gas Leak Sensor ML Rebuild | {index}"
    p.font.size = Pt(9)
    p.font.color.rgb = MID
    p.alignment = PP_ALIGN.RIGHT


def title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.35), prs.slide_width, Inches(1.15))
    set_fill(bar, TEAL)
    box = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(11.8), Inches(1.05))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = WHITE
    sub = slide.shapes.add_textbox(Inches(0.78), Inches(3.35), Inches(10.8), Inches(0.55))
    sp = sub.text_frame.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(20)
    sp.font.color.rgb = RGBColor(220, 235, 239)
    tag = slide.shapes.add_textbox(Inches(0.78), Inches(6.67), Inches(9.2), Inches(0.35))
    tp = tag.text_frame.paragraphs[0]
    tp.text = "ESP32-S3 + 8 MQ Sensor Array + TensorFlow Lite Micro + LoRa"
    tp.font.size = Pt(15)
    tp.font.bold = True
    tp.font.color.rgb = WHITE


def bullets_slide(title, bullets, subtitle=None, accent=TEAL):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, title, subtitle)
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.38), Inches(0.12), Inches(5.25))
    set_fill(stripe, accent)
    box = slide.shapes.add_textbox(Inches(0.95), Inches(1.35), Inches(11.5), Inches(5.35))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, text in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = NAVY
        p.space_after = Pt(10)
    add_footer(slide, len(prs.slides))


def two_column_slide(title, left_title, left_items, right_title, right_items, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, title, subtitle)
    add_panel(slide, Inches(0.65), Inches(1.55), Inches(5.85), Inches(5.25), left_title, left_items, TEAL)
    add_panel(slide, Inches(6.85), Inches(1.55), Inches(5.85), Inches(5.25), right_title, right_items, GREEN)
    add_footer(slide, len(prs.slides))


def add_panel(slide, x, y, w, h, heading, items, accent):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    set_fill(panel, WHITE)
    panel.line.color.rgb = RGBColor(220, 226, 231)
    panel.line.width = Pt(1)
    cap = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.48))
    set_fill(cap, accent)
    ht = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.11), w - Inches(0.5), Inches(0.28))
    p = ht.text_frame.paragraphs[0]
    p.text = heading
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE
    body = slide.shapes.add_textbox(x + Inches(0.28), y + Inches(0.75), w - Inches(0.55), h - Inches(1.0))
    tf = body.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(items):
        bp = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        bp.text = item
        bp.font.size = Pt(15)
        bp.font.color.rgb = NAVY
        bp.space_after = Pt(8)


def table_like_slide(title, rows, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, title, subtitle)
    y = Inches(1.45)
    row_h = Inches(0.55)
    for idx, (left, right) in enumerate(rows):
        color = WHITE if idx % 2 == 0 else RGBColor(237, 242, 245)
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), y, Inches(11.85), row_h)
        set_fill(rect, color)
        rect.line.color.rgb = RGBColor(220, 226, 231)
        lbox = slide.shapes.add_textbox(Inches(1.0), y + Inches(0.12), Inches(3.2), Inches(0.25))
        lp = lbox.text_frame.paragraphs[0]
        lp.text = left
        lp.font.size = Pt(13)
        lp.font.bold = True
        lp.font.color.rgb = NAVY
        rbox = slide.shapes.add_textbox(Inches(4.35), y + Inches(0.12), Inches(7.8), Inches(0.25))
        rp = rbox.text_frame.paragraphs[0]
        rp.text = right
        rp.font.size = Pt(13)
        rp.font.color.rgb = NAVY
        y += row_h
    add_footer(slide, len(prs.slides))


def image_review_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Gas Test Chamber Drawing Review", "Current mechanical concept and required revisions")
    if CHAMBER_IMAGE.exists():
        slide.shapes.add_picture(str(CHAMBER_IMAGE), Inches(0.55), Inches(1.45), width=Inches(6.15))
    add_panel(
        slide,
        Inches(7.0),
        Inches(1.45),
        Inches(5.65),
        Inches(5.35),
        "Revision Priorities",
        [
            "Add gas inlet, clean-air inlet, exhaust outlet, cable gland, and reference sampling port.",
            "Add gasket/seal detail around the removable plate.",
            "Add internal low-voltage mixing fan and log fan state.",
            "Replace internal wood with acrylic, aluminum, stainless, or PTFE.",
            "Consider reducing chamber volume from about 60 L to 10-30 L for early tests.",
            "Keep solenoids, relays, and switching electronics outside the gas space.",
        ],
        AMBER,
    )
    add_footer(slide, len(prs.slides))


title_slide("Gas Leak Sensor ML Rebuild", "Production-oriented roadmap for sensor-array gas detection")

bullets_slide(
    "Source Documents",
    [
        "design.md: system architecture, firmware contract, payload, risks, chamber design, and ML roadmap.",
        "documentation.md: client-facing explanation, model method, chamber plan, safety notes, and acceptance criteria.",
        "session_notes.md: current continuation notes and decisions from the rebuild session.",
        "todo.md: execution checklist for data, chamber, ML, firmware, gateway, and documentation work.",
        "Gas Test Chamber Design.jpeg: mechanical chamber drawing reviewed for revision recommendations.",
    ],
    "This deck summarizes the latest project documentation, not the older three-class prototype presentation.",
)

bullets_slide(
    "Project Goal",
    [
        "Detect gas leaks using the combined response fingerprint of 8 MQ sensors.",
        "Train one model per board because MQ boards behave differently.",
        "Run inference locally on ESP32-S3 using TensorFlow Lite Micro.",
        "Send versioned binary results through LoRa to the cluster head and gateway.",
        "Use chamber data and calibrated references before claiming production ppm.",
    ],
)

two_column_slide(
    "System Architecture",
    "Sensor Node",
    [
        "8 MQ voltages through ADS1256",
        "Board-specific scaler",
        "TFLite Micro inference",
        "Local alarm decision",
        "LoRa payload transmit",
    ],
    "Network And Reporting",
    [
        "Cluster head forwards LoRa packets",
        "Gateway publishes upstream data",
        "Future gateway decodes named MQTT fields",
        "Payload stays under 64-byte transport limit",
        "Model/scaler/payload versions are tracked",
    ],
)

table_like_slide(
    "Gas Class Contract",
    [
        ("0", "normal"),
        ("1", "methane"),
        ("2", "H2S"),
        ("3", "butane / LPG-related gas"),
        ("4", "propane"),
        ("5", "CO"),
    ],
    "Stable IDs prevent firmware and dashboard churn as new chamber data is added.",
)

bullets_slide(
    "Current Dataset Status",
    [
        "Usable Excel datasets exist for Board1, Board3-Board7, Board9, Board10, and Board11.",
        "Board2-2 has firmware but no matching Excel dataset; Board8 dataset is empty.",
        "Current data mostly covers normal, methane, and butane/LPG-related rows.",
        "No calibrated ppm labels exist yet, so ppm output is only a response proxy.",
        "The dataset is useful for prototype validation, not production safety claims.",
    ],
    accent=AMBER,
)

bullets_slide(
    "Production Dataset Requirements",
    [
        "Collect time-series chamber data, not isolated rows.",
        "Log session_id, board_id, gas_type, param, reference_ppm, BME280 readings, and chamber state.",
        "Capture baseline, injection, mixing/rise, stable target ppm, recovery, and post-recovery phases.",
        "Include false-positive vapors such as alcohol, perfume, smoke, cleaning chemicals, humid air, and exhaust-like cases.",
        "Validate on held-out sessions or days instead of random row splits.",
    ],
)

two_column_slide(
    "Model Design",
    "Current Improvement Program",
    [
        "Input: 8 MQ voltages",
        "Shared dense body",
        "gas_type softmax",
        "leak_present sigmoid",
        "severity softmax",
        "ppm_estimate regression proxy",
    ],
    "Next Production Model",
    [
        "Add BME280 features",
        "Add baseline delta and ratio features",
        "Add 10-60 second window features",
        "Add slope, moving average, max response",
        "Add uncertain / out-of-distribution behavior",
    ],
)

bullets_slide(
    "Suggested ESP32-S3 Model Shape",
    [
        "Inputs: 8 scaled MQ voltages + 8 baseline deltas + 8 slope/window features + temperature + humidity + pressure.",
        "Network: Dense 64 ReLU -> Dense 32 ReLU -> Shared Dense 16 ReLU.",
        "Outputs: six-class gas_type, leak_present, four-class severity, ppm_estimate.",
        "Keep float TFLite export until firmware explicitly supports int8 quantized tensors.",
        "Measure tensor arena size and inference time on device before expanding the model.",
    ],
)

bullets_slide(
    "Validation Metrics",
    [
        "Report precision, recall, F1, and confusion matrix per gas, not only total accuracy.",
        "Track false alarm rate and missed leak rate; missed leak rate is the key safety metric.",
        "Report leak-present precision/recall and severity confusion matrix.",
        "Report ppm MAE per gas only after calibrated reference_ppm labels exist.",
        "Break metrics down by board, session/day, humidity range, and temperature range.",
    ],
    accent=GREEN,
)

bullets_slide(
    "Firmware Contract",
    [
        "Verify ADC channel order exactly matches MQ135V, MQ2V, MQ3V, MQ4V, MQ7V, MQ5V, MQ6V, MQ8V.",
        "Validate input tensor type, input shape, output count, and output dimensions before trusting inference.",
        "Fail safe on invalid voltages, disconnected sensors, stale readings, ADC saturation, low confidence, or OOD inputs.",
        "Include model version, scaler version, payload version, and board ID in artifacts and reports.",
        "Keep alarm threshold conservative until session/day validation proves reliability.",
    ],
    accent=RED,
)

table_like_slide(
    "Versioned LoRa Payload",
    [
        ("Size", "32 bytes, fits current 64-byte message limit"),
        ("Header", "version, gas_type, leak_present, severity"),
        ("Confidence", "gas, leak, and severity confidence scaled x1000"),
        ("Estimate", "uint16 ppm proxy now; real ppm later after calibration"),
        ("Timing", "uint32 inference time in microseconds"),
        ("Sensors", "8 signed int16 MQ millivolt readings"),
    ],
)

bullets_slide(
    "Chamber Hardware Recommendation",
    [
        "Use a controlled-flow chamber, not a sealed box with manual gas injection.",
        "Recommended early chamber size: 10-30 L clear acrylic or polycarbonate with gasketed lid.",
        "Add clean-air inlet, gas inlet, exhaust outlet, cable glands, mixing fan, and safe ventilation.",
        "Keep solenoids, relays, and switching contacts outside the gas space.",
        "Use regulator plus needle valve, flow meter, or mass flow controller for repeatable tests.",
    ],
)

image_review_slide()

table_like_slide(
    "Reference Instruments",
    [
        ("Butane / LPG", "Calibrated LPG or LEL detector; TGS2610 as reference/cross-check"),
        ("Propane", "Calibrated LPG or LEL detector; TGS2610 as reference/cross-check"),
        ("Methane", "Calibrated methane or LEL detector"),
        ("CO", "Calibrated CO detector/logger or CO analyzer"),
        ("H2S", "Calibrated H2S detector/logger or H2S analyzer"),
    ],
    "CO and H2S require certified instruments, ventilation, gas-rated safety equipment, and supervision.",
)

bullets_slide(
    "Critical Risks",
    [
        "Random row split can overestimate accuracy because adjacent chamber samples are correlated.",
        "Current ppm_estimate is not real calibrated concentration.",
        "Missing H2S, propane, and CO data blocks production claims for those gases.",
        "Swapped ADC feature order can silently break on-device inference.",
        "False-positive vapors and long-term drift must be tested before deployment claims.",
    ],
    accent=RED,
)

bullets_slide(
    "Next Work",
    [
        "Finalize chamber dataset schema and operating procedure.",
        "Revise chamber drawing with controlled-flow ports, gasket, fan, cable glands, and safer internal materials.",
        "Build chamber controller/logger for MQ, reference sensor, BME280, valve states, fan state, and chamber state.",
        "Collect LPG-family time-series data first, then expand to methane with calibrated reference hardware.",
        "Integrate the improved four-output model wrapper into firmware after tensor contract validation.",
    ],
    accent=TEAL,
)

prs.save(OUT)
print(f"Presentation saved: {OUT}")
