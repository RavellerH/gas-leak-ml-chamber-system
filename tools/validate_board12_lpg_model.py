"""Validation report for Board12 clean-air vs LPG detector.

Outputs Markdown, Excel, JSON metrics, plots, and a PowerPoint deck under
test LPG board 12/ml_validation.
"""

from __future__ import annotations

import json
import math
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from sklearn.dummy import DummyClassifier
from sklearn.linear_model import LogisticRegression
from sklearn.metrics import (
    accuracy_score,
    balanced_accuracy_score,
    classification_report,
    confusion_matrix,
    f1_score,
    precision_score,
    recall_score,
    roc_auc_score,
)
from sklearn.model_selection import GroupKFold, StratifiedKFold, train_test_split
from sklearn.pipeline import Pipeline
from sklearn.preprocessing import StandardScaler

from train_board12_lpg_detector import FEATURE_COLS, MODEL_FEATURES, add_features, auto_label, load_raw_csvs


DATA_DIR = Path("test LPG board 12")
OUT_DIR = DATA_DIR / "ml_validation"
MODEL_JSON = DATA_DIR / "ml_board12_lpg_detector" / "board12_lpg_detector_model.json"
LABELS = ["clean_air", "lpg_detected"]


def y_from_labels(labels: pd.Series) -> np.ndarray:
    return (labels == "lpg_detected").astype(int).to_numpy()


def make_model() -> Pipeline:
    return Pipeline(
        [
            ("scaler", StandardScaler()),
            ("model", LogisticRegression(class_weight="balanced", max_iter=1000, random_state=42)),
        ]
    )


def summarize_predictions(y_true, y_pred, y_prob=None) -> dict:
    metrics = {
        "accuracy": float(accuracy_score(y_true, y_pred)),
        "balanced_accuracy": float(balanced_accuracy_score(y_true, y_pred)),
        "precision_lpg": float(precision_score(y_true, y_pred, zero_division=0)),
        "recall_lpg": float(recall_score(y_true, y_pred, zero_division=0)),
        "f1_lpg": float(f1_score(y_true, y_pred, zero_division=0)),
        "confusion_matrix": confusion_matrix(y_true, y_pred, labels=[0, 1]).tolist(),
        "support_clean_air": int((np.asarray(y_true) == 0).sum()),
        "support_lpg_detected": int((np.asarray(y_true) == 1).sum()),
    }
    if y_prob is not None and len(set(y_true)) == 2:
        metrics["roc_auc"] = float(roc_auc_score(y_true, y_prob))
    return metrics


def threshold_predict(train_x: pd.DataFrame, train_y: np.ndarray, test_x: pd.DataFrame, feature: str) -> tuple[np.ndarray, float]:
    clean = train_x.loc[train_y == 0, feature]
    lpg = train_x.loc[train_y == 1, feature]
    threshold = float((clean.quantile(0.95) + lpg.quantile(0.05)) / 2.0)
    return (test_x[feature].to_numpy() >= threshold).astype(int), threshold


def response_score_predict(train_x: pd.DataFrame, train_y: np.ndarray, test_x: pd.DataFrame) -> tuple[np.ndarray, float]:
    clean = train_x.loc[train_y == 0, FEATURE_COLS]
    baseline = clean.median()
    spread = clean.std(ddof=0).clip(lower=0.0002)

    def score(x: pd.DataFrame) -> pd.Series:
        return ((x[FEATURE_COLS] - baseline).abs() / spread).mean(axis=1)

    train_score = score(train_x)
    test_score = score(test_x)
    clean_score = train_score[train_y == 0]
    lpg_score = train_score[train_y == 1]
    threshold = float((clean_score.quantile(0.95) + lpg_score.quantile(0.05)) / 2.0)
    return (test_score.to_numpy() >= threshold).astype(int), threshold


def random_split_eval(df: pd.DataFrame) -> dict:
    x = df[MODEL_FEATURES]
    y = y_from_labels(df["auto_label"])
    train_idx, test_idx = train_test_split(np.arange(len(df)), test_size=0.2, random_state=42, stratify=y)
    x_train, x_test = x.iloc[train_idx], x.iloc[test_idx]
    y_train, y_test = y[train_idx], y[test_idx]

    model = make_model()
    model.fit(x_train, y_train)
    train_pred = model.predict(x_train)
    test_pred = model.predict(x_test)
    test_prob = model.predict_proba(x_test)[:, 1]

    baselines = {}
    for feature in ["mean_v", "max_v"]:
        pred, threshold = threshold_predict(x_train, y_train, x_test, feature)
        baselines[f"{feature}_threshold"] = {"threshold": threshold, **summarize_predictions(y_test, pred)}
    pred, threshold = response_score_predict(x_train, y_train, x_test)
    baselines["response_score_threshold"] = {"threshold": threshold, **summarize_predictions(y_test, pred)}

    dummy = DummyClassifier(strategy="most_frequent")
    dummy.fit(x_train, y_train)
    dummy_pred = dummy.predict(x_test)
    baselines["majority_class"] = summarize_predictions(y_test, dummy_pred)

    return {
        "train": summarize_predictions(y_train, train_pred),
        "test": summarize_predictions(y_test, test_pred, test_prob),
        "baselines": baselines,
        "test_predictions": pd.DataFrame(
            {
                "row_index": test_idx,
                "source_file": df.iloc[test_idx]["source_file"].to_numpy(),
                "actual": y_test,
                "predicted": test_pred,
                "lpg_probability": test_prob,
                "mean_v": x_test["mean_v"].to_numpy(),
                "max_v": x_test["max_v"].to_numpy(),
            }
        ),
    }


def stratified_kfold_eval(df: pd.DataFrame, folds: int = 5) -> tuple[dict, pd.DataFrame]:
    x = df[MODEL_FEATURES]
    y = y_from_labels(df["auto_label"])
    skf = StratifiedKFold(n_splits=folds, shuffle=True, random_state=42)
    rows = []
    for fold, (train_idx, test_idx) in enumerate(skf.split(x, y), start=1):
        model = make_model()
        model.fit(x.iloc[train_idx], y[train_idx])
        pred = model.predict(x.iloc[test_idx])
        prob = model.predict_proba(x.iloc[test_idx])[:, 1]
        m = summarize_predictions(y[test_idx], pred, prob)
        rows.append({"fold": fold, **{k: v for k, v in m.items() if k != "confusion_matrix"}})
    folds_df = pd.DataFrame(rows)
    summary = {
        f"{col}_mean": float(folds_df[col].mean())
        for col in folds_df.columns
        if col != "fold"
    }
    summary.update(
        {
            f"{col}_std": float(folds_df[col].std(ddof=0))
            for col in folds_df.columns
            if col != "fold"
        }
    )
    return summary, folds_df


def group_holdout_eval(df: pd.DataFrame) -> tuple[dict, pd.DataFrame]:
    x = df[MODEL_FEATURES]
    y = y_from_labels(df["auto_label"])
    groups = df["source_file"].to_numpy()
    rows = []
    predictions = []
    for heldout in sorted(df["source_file"].unique()):
        train_idx = np.flatnonzero(groups != heldout)
        test_idx = np.flatnonzero(groups == heldout)
        if len(set(y[train_idx])) < 2 or len(set(y[test_idx])) < 2:
            note = "single-class train/test split; metrics may be limited"
        else:
            note = ""
        model = make_model()
        model.fit(x.iloc[train_idx], y[train_idx])
        pred = model.predict(x.iloc[test_idx])
        prob = model.predict_proba(x.iloc[test_idx])[:, 1]
        m = summarize_predictions(y[test_idx], pred, prob if len(set(y[test_idx])) == 2 else None)
        rows.append({"heldout_file": heldout, "note": note, **{k: v for k, v in m.items() if k != "confusion_matrix"}})
        predictions.append(
            pd.DataFrame(
                {
                    "heldout_file": heldout,
                    "row_index": test_idx,
                    "actual": y[test_idx],
                    "predicted": pred,
                    "lpg_probability": prob,
                    "mean_v": x.iloc[test_idx]["mean_v"].to_numpy(),
                }
            )
        )
    result = pd.DataFrame(rows)
    summary = {
        "group_holdout_accuracy_mean": float(result["accuracy"].mean()),
        "group_holdout_accuracy_min": float(result["accuracy"].min()),
        "group_holdout_f1_lpg_mean": float(result["f1_lpg"].mean()),
        "group_holdout_f1_lpg_min": float(result["f1_lpg"].min()),
    }
    return summary, result, pd.concat(predictions, ignore_index=True)


def group_kfold_eval(df: pd.DataFrame) -> tuple[dict, pd.DataFrame]:
    x = df[MODEL_FEATURES]
    y = y_from_labels(df["auto_label"])
    groups = df["source_file"].to_numpy()
    n_groups = len(np.unique(groups))
    n_splits = min(5, n_groups)
    gkf = GroupKFold(n_splits=n_splits)
    rows = []
    for fold, (train_idx, test_idx) in enumerate(gkf.split(x, y, groups), start=1):
        model = make_model()
        model.fit(x.iloc[train_idx], y[train_idx])
        pred = model.predict(x.iloc[test_idx])
        prob = model.predict_proba(x.iloc[test_idx])[:, 1]
        m = summarize_predictions(y[test_idx], pred, prob if len(set(y[test_idx])) == 2 else None)
        rows.append(
            {
                "fold": fold,
                "heldout_files": ", ".join(sorted(set(groups[test_idx]))),
                **{k: v for k, v in m.items() if k != "confusion_matrix"},
            }
        )
    folds_df = pd.DataFrame(rows)
    summary = {
        "group_kfold_splits": n_splits,
        "group_kfold_accuracy_mean": float(folds_df["accuracy"].mean()),
        "group_kfold_accuracy_std": float(folds_df["accuracy"].std(ddof=0)),
        "group_kfold_f1_lpg_mean": float(folds_df["f1_lpg"].mean()),
        "group_kfold_f1_lpg_std": float(folds_df["f1_lpg"].std(ddof=0)),
    }
    return summary, folds_df


def dataset_audit(all_df: pd.DataFrame, train_df: pd.DataFrame) -> tuple[pd.DataFrame, dict]:
    rows = []
    for name, g in all_df.groupby("source_file", sort=True):
        elapsed = g["elapsed_ms"]
        diffs = elapsed.sort_values().diff().dropna()
        rows.append(
            {
                "source_file": name,
                "rows": int(len(g)),
                "used_rows": int(g["auto_label"].isin(LABELS).sum()),
                "clean_air": int((g["auto_label"] == "clean_air").sum()),
                "lpg_detected": int((g["auto_label"] == "lpg_detected").sum()),
                "ambiguous": int((g["auto_label"] == "ambiguous").sum()),
                "duration_s": float((elapsed.max() - elapsed.min()) / 1000.0),
                "sample_interval_median_s": float(diffs.median() / 1000.0) if len(diffs) else 0.0,
                "sample_interval_max_s": float(diffs.max() / 1000.0) if len(diffs) else 0.0,
                "mean_v_min": float(g["mean_v"].min()),
                "mean_v_max": float(g["mean_v"].max()),
                "mean_v_first": float(g["mean_v"].iloc[0]),
                "mean_v_last": float(g["mean_v"].iloc[-1]),
            }
        )
    audit_df = pd.DataFrame(rows)
    missing = int(all_df[FEATURE_COLS + ["elapsed_ms"]].isna().sum().sum())
    duplicates = int(all_df.duplicated(subset=["source_file", "elapsed_ms", *FEATURE_COLS]).sum())
    summary = {
        "raw_rows": int(len(all_df)),
        "training_rows": int(len(train_df)),
        "ambiguous_rows_excluded": int((all_df["auto_label"] == "ambiguous").sum()),
        "missing_values_after_load": missing,
        "duplicate_rows": duplicates,
        "source_files": int(all_df["source_file"].nunique()),
        "class_counts_training": train_df["auto_label"].value_counts().to_dict(),
    }
    return audit_df, summary


def browser_parity_check(df: pd.DataFrame) -> tuple[dict, pd.DataFrame]:
    if not MODEL_JSON.exists():
        return {"available": False}, pd.DataFrame()
    model_json = json.loads(MODEL_JSON.read_text(encoding="utf-8"))
    sample = df.sample(n=min(200, len(df)), random_state=7).copy()
    x = sample[model_json["feature_cols"]]
    coef = np.asarray(model_json["coef"], dtype=float)
    mean = np.asarray(model_json["scaler_mean"], dtype=float)
    scale = np.asarray(model_json["scaler_scale"], dtype=float)
    z = ((x.to_numpy(dtype=float) - mean) / scale).dot(coef) + float(model_json["intercept"])
    prob_json = 1.0 / (1.0 + np.exp(-z))
    model = make_model()
    full_x = df[MODEL_FEATURES]
    full_y = y_from_labels(df["auto_label"])
    model.fit(full_x, full_y)
    prob_py = model.predict_proba(x)[:, 1]
    parity = pd.DataFrame(
        {
            "source_file": sample["source_file"].to_numpy(),
            "auto_label": sample["auto_label"].to_numpy(),
            "json_probability": prob_json,
            "retrained_python_probability": prob_py,
            "absolute_difference": np.abs(prob_json - prob_py),
        }
    )
    # The JSON was exported by an earlier training run; exact equality is not expected after retraining.
    summary = {
        "available": True,
        "checked_rows": int(len(parity)),
        "json_probability_min": float(parity["json_probability"].min()),
        "json_probability_max": float(parity["json_probability"].max()),
        "python_probability_min": float(parity["retrained_python_probability"].min()),
        "python_probability_max": float(parity["retrained_python_probability"].max()),
        "mean_absolute_difference": float(parity["absolute_difference"].mean()),
        "note": "Compares browser JSON artifact with a freshly retrained Python model on the same rows; not expected to be identical unless artifacts are regenerated.",
    }
    return summary, parity


def save_plots(all_df: pd.DataFrame, train_df: pd.DataFrame, random_preds: pd.DataFrame, group_holdout: pd.DataFrame) -> dict:
    paths = {}
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    plt.figure(figsize=(10, 5))
    for label, color in [("clean_air", "#177245"), ("lpg_detected", "#b42318"), ("ambiguous", "#a15c07")]:
        g = all_df[all_df["auto_label"] == label]
        plt.scatter(g["elapsed_s"], g["mean_v"], s=8, alpha=0.55, label=label, c=color)
    plt.xlabel("Elapsed time within source file (s)")
    plt.ylabel("Mean MQ voltage (V)")
    plt.title("Auto-labeled Board12 LPG Dataset")
    plt.legend()
    plt.grid(alpha=0.25)
    plt.tight_layout()
    paths["label_scatter"] = OUT_DIR / "label_scatter.png"
    plt.savefig(paths["label_scatter"], dpi=160)
    plt.close()

    plt.figure(figsize=(9, 5))
    for actual, color, name in [(0, "#177245", "clean_air"), (1, "#b42318", "lpg_detected")]:
        vals = random_preds[random_preds["actual"] == actual]["lpg_probability"]
        plt.hist(vals, bins=30, alpha=0.7, label=name, color=color)
    plt.axvline(0.8, color="#b42318", linestyle="--", label="80% LPG threshold")
    plt.xlabel("LPG probability")
    plt.ylabel("Rows")
    plt.title("Random 80/20 Probability Distribution")
    plt.legend()
    plt.tight_layout()
    paths["probability_hist"] = OUT_DIR / "probability_hist.png"
    plt.savefig(paths["probability_hist"], dpi=160)
    plt.close()

    plt.figure(figsize=(10, 5))
    sorted_group = group_holdout.sort_values("accuracy")
    plt.bar(sorted_group["heldout_file"], sorted_group["accuracy"], color="#0f766e")
    plt.xticks(rotation=35, ha="right")
    plt.ylim(0, 1.05)
    plt.ylabel("Accuracy")
    plt.title("Capture-file Holdout Accuracy")
    plt.tight_layout()
    paths["group_holdout_accuracy"] = OUT_DIR / "group_holdout_accuracy.png"
    plt.savefig(paths["group_holdout_accuracy"], dpi=160)
    plt.close()

    return {k: str(v) for k, v in paths.items()}


def write_markdown(summary: dict, audit_df: pd.DataFrame, random_result: dict, skf_summary: dict, group_summary: dict, gkf_summary: dict) -> None:
    lines = [
        "# Board12 LPG ML Validation Report",
        "",
        "Scope: Board12 only, binary clean_air vs lpg_detected. Ambiguous transition rows are excluded from training and reported separately.",
        "",
        "## Dataset Quality",
        f"- Raw rows loaded: {summary['dataset']['raw_rows']:,}",
        f"- Training/evaluation rows: {summary['dataset']['training_rows']:,}",
        f"- Ambiguous rows excluded: {summary['dataset']['ambiguous_rows_excluded']:,}",
        f"- Source files: {summary['dataset']['source_files']}",
        f"- Missing values after load: {summary['dataset']['missing_values_after_load']}",
        f"- Duplicate rows: {summary['dataset']['duplicate_rows']}",
        "",
        "## Labeling",
        f"- Clean reference mean: {summary['labeling']['clean_reference_mean_v']:.6f} V",
        f"- Clean limit: {summary['labeling']['clean_limit_v']:.6f} V mean",
        f"- LPG limit: {summary['labeling']['lpg_limit_v']:.6f} V mean",
        f"- Training class counts: {summary['dataset']['class_counts_training']}",
        "",
        "## Random 80/20 Split",
        f"- Train accuracy: {random_result['train']['accuracy']:.4f}",
        f"- Test accuracy: {random_result['test']['accuracy']:.4f}",
        f"- Test balanced accuracy: {random_result['test']['balanced_accuracy']:.4f}",
        f"- Test LPG precision / recall / F1: {random_result['test']['precision_lpg']:.4f} / {random_result['test']['recall_lpg']:.4f} / {random_result['test']['f1_lpg']:.4f}",
        "",
        "## Cross-validation",
        f"- Stratified 5-fold accuracy: {skf_summary['accuracy_mean']:.4f} +/- {skf_summary['accuracy_std']:.4f}",
        f"- Stratified 5-fold LPG F1: {skf_summary['f1_lpg_mean']:.4f} +/- {skf_summary['f1_lpg_std']:.4f}",
        f"- Group K-fold accuracy: {gkf_summary['group_kfold_accuracy_mean']:.4f} +/- {gkf_summary['group_kfold_accuracy_std']:.4f}",
        f"- Group K-fold LPG F1: {gkf_summary['group_kfold_f1_lpg_mean']:.4f} +/- {gkf_summary['group_kfold_f1_lpg_std']:.4f}",
        "",
        "## Capture-file Holdout",
        f"- Mean accuracy: {group_summary['group_holdout_accuracy_mean']:.4f}",
        f"- Worst accuracy: {group_summary['group_holdout_accuracy_min']:.4f}",
        f"- Mean LPG F1: {group_summary['group_holdout_f1_lpg_mean']:.4f}",
        f"- Worst LPG F1: {group_summary['group_holdout_f1_lpg_min']:.4f}",
        "",
        "## Baseline Comparisons",
    ]
    for name, metrics in random_result["baselines"].items():
        thresh = metrics.get("threshold")
        prefix = f"- {name}: accuracy {metrics['accuracy']:.4f}, LPG F1 {metrics['f1_lpg']:.4f}"
        if thresh is not None:
            prefix += f", threshold {thresh:.6f}"
        lines.append(prefix)

    train_acc = random_result["train"]["accuracy"]
    test_acc = random_result["test"]["accuracy"]
    gap = train_acc - test_acc
    lines.extend(
        [
            "",
            "## Overfitting Assessment",
            f"- Train-test accuracy gap: {gap:.4f}",
            "- Random split is optimistic because adjacent time samples are highly correlated.",
            "- File-level holdout and group K-fold are more important indicators for future runs.",
            "- If simple thresholds match ML performance, an embedded threshold detector may be preferable for firmware.",
            "",
            "## Recommendations",
            "- Keep binary clean_air vs lpg_detected for now.",
            "- Continue excluding ambiguous/rising/clearing rows until they are explicitly labeled.",
            "- Add confirmed clean-air-only runs after full ventilation to reduce false-positive risk.",
            "- Add repeated non-saturated LPG runs for stronger file-holdout validation.",
        ]
    )
    (OUT_DIR / "Board12_LPG_ML_Validation_Report.md").write_text("\n".join(lines), encoding="utf-8")


def write_excel(outputs: dict) -> None:
    with pd.ExcelWriter(OUT_DIR / "Board12_LPG_ML_Validation.xlsx", engine="openpyxl") as writer:
        outputs["audit_df"].to_excel(writer, sheet_name="dataset_audit", index=False)
        outputs["all_df"].to_excel(writer, sheet_name="all_labeled_rows", index=False)
        outputs["train_df"].to_excel(writer, sheet_name="training_rows", index=False)
        pd.DataFrame([outputs["random_result"]["train"]]).drop(columns=["confusion_matrix"]).to_excel(writer, sheet_name="random_train", index=False)
        pd.DataFrame([outputs["random_result"]["test"]]).drop(columns=["confusion_matrix"]).to_excel(writer, sheet_name="random_test", index=False)
        outputs["random_result"]["test_predictions"].to_excel(writer, sheet_name="random_test_predictions", index=False)
        pd.DataFrame(outputs["random_result"]["baselines"]).T.to_excel(writer, sheet_name="baseline_comparison")
        outputs["skf_folds"].to_excel(writer, sheet_name="stratified_kfold", index=False)
        outputs["group_holdout_df"].to_excel(writer, sheet_name="file_holdout", index=False)
        outputs["group_kfold_df"].to_excel(writer, sheet_name="group_kfold", index=False)
        outputs["parity_df"].to_excel(writer, sheet_name="browser_parity", index=False)


def write_powerpoint(summary: dict, plot_paths: dict) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def title(slide, text, sub=""):
        box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12), Inches(0.6))
        p = box.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(30)
        p.font.bold = True
        if sub:
            s = slide.shapes.add_textbox(Inches(0.58), Inches(0.98), Inches(12), Inches(0.35))
            sp = s.text_frame.paragraphs[0]
            sp.text = sub
            sp.font.size = Pt(13)

    def bullets(slide, items, x=0.8, y=1.4, w=11.7, h=5.4, size=18):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.font.size = Pt(size)
            p.space_after = Pt(8)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title(slide, "Board12 LPG ML Validation", "Clean air vs LPG detected, ambiguous rows excluded")
    bullets(
        slide,
        [
            f"Raw rows: {summary['dataset']['raw_rows']:,}",
            f"Training/evaluation rows: {summary['dataset']['training_rows']:,}",
            f"Ambiguous rows excluded: {summary['dataset']['ambiguous_rows_excluded']:,}",
            f"Source files: {summary['dataset']['source_files']}",
            "Validation includes random 80/20 split, stratified k-fold, group/file holdout, baseline thresholds, and browser artifact parity.",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title(slide, "Dataset Labels and Quality")
    if "label_scatter" in plot_paths:
        slide.shapes.add_picture(plot_paths["label_scatter"], Inches(0.65), Inches(1.25), width=Inches(7.3))
    bullets(
        slide,
        [
            f"Clean limit: {summary['labeling']['clean_limit_v']:.6f} V mean",
            f"LPG limit: {summary['labeling']['lpg_limit_v']:.6f} V mean",
            f"Missing values after load: {summary['dataset']['missing_values_after_load']}",
            f"Duplicate rows: {summary['dataset']['duplicate_rows']}",
        ],
        x=8.25,
        y=1.45,
        w=4.4,
        h=4.8,
        size=16,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title(slide, "Random 80/20 and Probability Separation")
    if "probability_hist" in plot_paths:
        slide.shapes.add_picture(plot_paths["probability_hist"], Inches(0.75), Inches(1.25), width=Inches(7.2))
    bullets(
        slide,
        [
            f"Test accuracy: {summary['random_80_20']['test']['accuracy']:.4f}",
            f"LPG precision: {summary['random_80_20']['test']['precision_lpg']:.4f}",
            f"LPG recall: {summary['random_80_20']['test']['recall_lpg']:.4f}",
            "Random split can be optimistic because time-adjacent rows are similar.",
        ],
        x=8.25,
        y=1.45,
        w=4.4,
        h=4.8,
        size=16,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title(slide, "File-level Generalization")
    if "group_holdout_accuracy" in plot_paths:
        slide.shapes.add_picture(plot_paths["group_holdout_accuracy"], Inches(0.75), Inches(1.25), width=Inches(7.2))
    bullets(
        slide,
        [
            f"File-holdout mean accuracy: {summary['group_holdout']['group_holdout_accuracy_mean']:.4f}",
            f"File-holdout worst accuracy: {summary['group_holdout']['group_holdout_accuracy_min']:.4f}",
            f"Group K-fold accuracy: {summary['group_kfold']['group_kfold_accuracy_mean']:.4f} +/- {summary['group_kfold']['group_kfold_accuracy_std']:.4f}",
            "This is the key validation view for future test runs.",
        ],
        x=8.25,
        y=1.45,
        w=4.4,
        h=4.8,
        size=16,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title(slide, "Conclusions and Next Actions")
    bullets(
        slide,
        [
            "Binary clean_air vs LPG_detected is the right first target.",
            "Ambiguous transition rows should remain excluded until explicitly labeled.",
            "Compare ML against mean/max/response-score thresholds before embedding.",
            "Collect confirmed clean-air-only runs after full ventilation to reduce false positives.",
            "Collect repeated non-saturated LPG runs for stronger file-holdout validation.",
        ],
    )

    prs.save(OUT_DIR / "Board12_LPG_ML_Validation.pptx")


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    raw_df = load_raw_csvs(DATA_DIR)
    all_df, labeling = auto_label(raw_df)
    train_df = all_df[all_df["auto_label"].isin(LABELS)].copy()

    audit_df, dataset_summary = dataset_audit(all_df, train_df)
    random_result = random_split_eval(train_df)
    skf_summary, skf_folds = stratified_kfold_eval(train_df)
    group_summary, group_holdout_df, group_predictions = group_holdout_eval(train_df)
    gkf_summary, group_kfold_df = group_kfold_eval(train_df)
    parity_summary, parity_df = browser_parity_check(train_df)
    plot_paths = save_plots(all_df, train_df, random_result["test_predictions"], group_holdout_df)

    summary = {
        "dataset": dataset_summary,
        "labeling": labeling,
        "random_80_20": {
            "train": random_result["train"],
            "test": random_result["test"],
            "baselines": random_result["baselines"],
        },
        "stratified_kfold": skf_summary,
        "group_holdout": group_summary,
        "group_kfold": gkf_summary,
        "browser_parity": parity_summary,
        "plots": plot_paths,
    }

    (OUT_DIR / "validation_summary.json").write_text(json.dumps(summary, indent=2), encoding="utf-8")
    all_df.to_csv(OUT_DIR / "Board12_lpg_validation_all_labeled_rows.csv", index=False)
    train_df.to_csv(OUT_DIR / "Board12_lpg_validation_training_rows.csv", index=False)
    audit_df.to_csv(OUT_DIR / "dataset_audit.csv", index=False)
    group_holdout_df.to_csv(OUT_DIR / "file_holdout_metrics.csv", index=False)
    group_predictions.to_csv(OUT_DIR / "file_holdout_predictions.csv", index=False)
    skf_folds.to_csv(OUT_DIR / "stratified_kfold_metrics.csv", index=False)
    group_kfold_df.to_csv(OUT_DIR / "group_kfold_metrics.csv", index=False)
    parity_df.to_csv(OUT_DIR / "browser_parity_check.csv", index=False)

    outputs = {
        "all_df": all_df,
        "train_df": train_df,
        "audit_df": audit_df,
        "random_result": random_result,
        "skf_folds": skf_folds,
        "group_holdout_df": group_holdout_df,
        "group_kfold_df": group_kfold_df,
        "parity_df": parity_df,
    }
    write_excel(outputs)
    write_markdown(summary, audit_df, random_result, skf_summary, group_summary, gkf_summary)
    write_powerpoint(summary, plot_paths)

    print(json.dumps(summary, indent=2))
    print(f"Wrote validation outputs to {OUT_DIR}")


if __name__ == "__main__":
    main()
